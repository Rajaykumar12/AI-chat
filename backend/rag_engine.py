import os
from pypdf import PdfReader
from pptx import Presentation
from typing import List, Tuple

# Simple in-memory vector store (alternative to ChromaDB)
# This uses basic TF-IDF similarity for document retrieval
document_store: List[Tuple[str, str, dict]] = []  # (id, text, metadata)

def get_or_create_collection():
    """Compatibility function - returns document store."""
    return document_store

def extract_pdf_text(filepath):
    """Extract text from PDF file."""
    reader = PdfReader(filepath)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_pptx_text(filepath):
    """Extract text from PPTX file."""
    prs = Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def ingest_documents(folder_path="documents"):
    """Reads PDFs and PPTX files and stores them in memory."""
    global document_store
    
    # Check if already indexed
    if len(document_store) > 0:
        print(f"Documents already indexed ({len(document_store)} chunks).")
        return
    
    # Create folder if doesn't exist
    if not os.path.exists(folder_path):
        os.makedirs(folder_path)
        print(f"Created {folder_path} folder. Please add PDF/PPTX files.")
        return
    
    print("Ingesting documents...")
    
    for filename in os.listdir(folder_path):
        filepath = os.path.join(folder_path, filename)
        text = ""
        
        if filename.endswith(".pdf"):
            text = extract_pdf_text(filepath)
        elif filename.endswith(".pptx"):
            text = extract_pptx_text(filepath)
        else:
            continue
        
        # Chunk text (1000 chars with 200 char overlap)
        chunk_size = 1000
        overlap = 200
        chunks = []
        for i in range(0, len(text), chunk_size - overlap):
            chunk = text[i:i + chunk_size].strip()
            if chunk:
                chunks.append(chunk)
        
        # Add chunks to document store
        for idx, chunk in enumerate(chunks):
            doc_id = f"{filename}_{idx}"
            metadata = {"source": filename, "chunk_id": idx}
            document_store.append((doc_id, chunk, metadata))
    
    if document_store:
        unique_files = len(set([m[2]['source'] for m in document_store]))
        print(f"Indexed {len(document_store)} chunks from {unique_files} documents.")
    else:
        print(f"No documents found in {folder_path}")

def simple_similarity(query: str, text: str) -> float:
    """Simple word overlap similarity (alternative to embeddings)."""
    query_words = set(query.lower().split())
    text_words = set(text.lower().split())
    
    if not query_words:
        return 0.0
    
    overlap = len(query_words.intersection(text_words))
    return overlap / len(query_words)

def retrieve_context(query: str, n_results=5):
    """Queries the document store for relevant content."""
    if len(document_store) == 0:
        return None
    
    # Score all documents
    scored_docs = []
    for doc_id, text, metadata in document_store:
        score = simple_similarity(query, text)
        scored_docs.append((score, text, metadata))
    
    # Sort by score and get top results
    scored_docs.sort(key=lambda x: x[0], reverse=True)
    top_docs = scored_docs[:n_results]
    
    # Filter out docs with score 0
    top_docs = [doc for doc in top_docs if doc[0] > 0]
    
    if not top_docs:
        return None
    
    # Combine retrieved chunks with source information
    context_parts = []
    for score, text, metadata in top_docs:
        context_parts.append(f"[Source: {metadata['source']}]\n{text}")
    
    return "\n\n---\n\n".join(context_parts)